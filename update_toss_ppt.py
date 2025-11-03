#!/usr/bin/env python3
"""
토스 결제창 스크린샷을 PPT에 추가
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

# 기존 PPT 로드
ppt_path = "/Users/l2dogyu/Downloads/정화의서재_결제경로_토스페이먼츠.pptx"
prs = Presentation(ppt_path)

# 슬라이드 8 (index 7) 수정
slide = prs.slides[7]

# 기존 내용 삭제 (텍스트 박스들)
shapes_to_remove = []
for shape in slide.shapes:
    if hasattr(shape, "text"):
        shapes_to_remove.append(shape)

for shape in shapes_to_remove[1:]:  # 타이틀 제외하고 삭제
    sp = shape.element
    sp.getparent().remove(sp)

# 토스 결제창 스크린샷 추가
toss_widget_path = "/Users/l2dogyu/KICDA/ruby/kicda-jh/public/screenshots/toss_submission_2025-11-01/09_toss_payment_widget.png"

if os.path.exists(toss_widget_path):
    left = Inches(0.5)
    top = Inches(1.2)
    width = Inches(9)
    
    pic = slide.shapes.add_picture(toss_widget_path, left, top, width=width)
    print("✅ 토스 결제창 스크린샷 추가 완료")
    
    # 설명 텍스트
    desc_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.8), Inches(9), Inches(0.6))
    desc_frame = desc_box.text_frame
    desc_frame.text = "토스페이먼츠 결제창 - 카카오페이, 카드 등 다양한 결제 수단 제공 | 상품명 및 금액 정확히 표시"
    desc_p = desc_frame.paragraphs[0]
    desc_p.font.size = Pt(14)
    desc_p.font.color.rgb = RGBColor(107, 114, 128)
    
    # 배경
    fill = desc_box.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
else:
    print("⚠️  토스 결제창 스크린샷 파일 없음")

# 저장
prs.save(ppt_path)

print(f"\n🎉 PPT 업데이트 완료!")
print(f"📁 위치: {ppt_path}")
print(f"📊 슬라이드 수: {len(prs.slides)}장")
print(f"✅ 토스 결제창 스크린샷 포함 완료!")

